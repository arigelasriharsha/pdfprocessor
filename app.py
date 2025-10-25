import os
import re
import nltk
import pandas as pd
import pdfplumber
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, pipeline
import fitz  # PyMuPDF for extracting images
from flask import Flask, request, jsonify, send_file,render_template
from flask_cors import CORS
import zipfile
import fitz  # PyMuPDF for extracting images
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from nltk.tokenize import sent_tokenize
from pptx.util import Inches, Pt
import PyPDF2


app = Flask(__name__)
CORS(app)
#nltk.data.path.append('C:/Users/allad/AppData/Local/Programs/Python/Python310/nltk_data')
#nltk.download('punkt_tab')

# Initialize directories
UPLOAD_DIR = 'uploads'
OUTPUT_DIR = 'outputs'
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, pipeline

model_name = "allenai/led-large-16384-arxiv"
print("Loading models...")

# Load model and tokenizer directly from Hugging Face
model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
tokenizer = AutoTokenizer.from_pretrained(model_name)

# Example summarization pipeline (optional)
summarization_pipeline = pipeline("summarization", model=model, tokenizer=tokenizer)


# Utility Functions
def preprocess_text(text):
    """
    Preprocess text to clean and format according to the requirements:
    1. Split sentences based on full stops.
    2. Add newlines for numbered lists or points.
    3. Ensure paragraph separation.
    4. Place headings/subheadings on new lines.
    """
    # Split sentences based on full stops or question marks
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)\s(?=\S)', text)

    # Add newlines for numbered lists or points
    formatted_text = ""
    for sentence in sentences:
        if re.match(r'^\d+\.|\d+\)', sentence.strip()):  # Numbered patterns like 1. or 1)
            formatted_text += f"\n{sentence.strip()}"
        else:
            formatted_text += f" {sentence.strip()}"

    # Split paragraphs based on newlines and ensure paragraph separation
    paragraphs = formatted_text.split("\n")
    formatted_paragraphs = []
    for paragraph in paragraphs:
        if paragraph.strip():
            # Check for headings/subheadings (e.g., uppercase or followed by a colon)
            if paragraph.strip().isupper() or paragraph.strip().endswith(":"):
                formatted_paragraphs.append(f"\n{paragraph.strip()}\n")
            else:
                formatted_paragraphs.append(paragraph.strip())

    return "\n\n".join(formatted_paragraphs)

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file while meeting the formatting requirements:
    - Avoid mixing table content with main text.
    - Return formatted text.
    """
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                # Extract plain text from the page
                page_text = page.extract_text() or ""

                # Avoid adding table content to plain text
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        table_content = " ".join([" ".join(row) for row in table if row])
                        page_text = page_text.replace(table_content, "")

                text += page_text

        return preprocess_text(text)

    except Exception as e:
        raise Exception(f"Error while extracting text: {e}")



def extract_tables_from_pdf(pdf_path):
    tables = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                extracted_tables = page.extract_tables()
                if extracted_tables:
                    tables.extend(extracted_tables)
        return tables
    except Exception as e:
        raise e

def save_tables_to_excel(tables, output_path):
    try:
        with pd.ExcelWriter(output_path) as writer:
            for i, table in enumerate(tables):
                df = pd.DataFrame(table)
                df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
        return output_path
    except Exception as e:
        raise e

# Chunk text by tokens
def chunk_text_by_tokens(text, max_tokens=16000):
    sentences = sent_tokenize(text)
    chunks, current_chunk = [], []
    current_chunk_tokens = 0
    for sentence in sentences:
        sentence_tokens = tokenizer.encode(sentence, add_special_tokens=False)
        sentence_token_count = len(sentence_tokens)
        if current_chunk_tokens + sentence_token_count > max_tokens:
            chunks.append(" ".join(current_chunk).strip())
            current_chunk = [sentence]
            current_chunk_tokens = sentence_token_count
        else:
            current_chunk.append(sentence)
            current_chunk_tokens += sentence_token_count
    if current_chunk:
        chunks.append(" ".join(current_chunk).strip())
    return chunks


def extract_important_points(text, num_points=5):
    """
    Fast, deterministic extraction of important points without calling large ML models.
    - Splits text into sentences using a lightweight regex.
    - Scores sentences by keyword-frequency (Counter) and position.
    - Returns path to a saved text file with up to `num_points` human-friendly bullets.
    """
    try:
        import os, re
        from collections import Counter

        # Basic validation & cleaning
        if not text or not isinstance(text, str) or not text.strip():
            raise ValueError("Input text is empty.")

        txt = re.sub(r'\s+', ' ', text).strip()

        # Lightweight sentence splitter (works reasonably well for English)
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', txt) if s.strip()]

        if not sentences:
            raise ValueError("No sentences found in the input text.")

        # Build keyword frequencies (words of length >=4)
        words = re.findall(r'\b[a-zA-Z]{4,}\b', txt.lower())
        if not words:
            # fallback: use whole sentences if no words found
            selection = sentences[:num_points]
            out = [f"• {s}" for s in selection]
            os.makedirs("outputs", exist_ok=True)
            out_path = os.path.join("outputs", "important_points.txt")
            with open(out_path, "w", encoding="utf-8") as f:
                f.write("\n".join(out))
            return out_path

        freq = Counter(words)
        top_keywords = set([w for w, _ in freq.most_common(40)])  # top tokens to consider

        # Score each sentence:
        #  - keyword score = sum of keyword frequencies that appear in the sentence
        #  - position bonus: earlier sentences get a small bonus
        scored = []
        for idx, s in enumerate(sentences):
            s_words = re.findall(r'\b[a-zA-Z]{4,}\b', s.lower())
            kw_score = sum(freq[w] for w in s_words if w in top_keywords)
            pos_bonus = max(0, (len(sentences) - idx) / len(sentences)) * 0.2  # earlier sentences slightly favored
            length_penalty = 0
            if len(s) < 40:   # too short -> small penalty
                length_penalty = 0.5
            if len(s) > 300:  # too long -> small penalty
                length_penalty += 0.3
            total_score = kw_score + pos_bonus - length_penalty
            scored.append((total_score, idx, s))

        # Sort by score (highest first), but keep original order for ties
        scored.sort(reverse=True, key=lambda x: (x[0], -x[1]))

        # Select candidate sentences, then re-order them by their original position for readability
        selected = []
        for score, idx, s in scored:
            # skip sentences that are obviously non-informative
            if len(re.sub(r'[^A-Za-z0-9]', '', s)) < 10:
                continue
            selected.append((idx, s))
            if len(selected) >= num_points * 3:  # gather a small pool then pick top-n by position
                break

        if not selected:
            # fallback to first N sentences
            selected = [(i, s) for i, s in enumerate(sentences[:num_points])]

        # Order by original position and pick up to num_points
        selected.sort(key=lambda x: x[0])
        final = [s for _, s in selected][:num_points]

        # Make points slightly more natural: ensure capitalization, end with period
        def tidy(sent):
            s = sent.strip()
            if not s:
                return ""
            s = s[0].upper() + s[1:]
            if s[-1] not in '.!?':
                s = s.rstrip() + '.'
            return s

        bullets = [f"• {tidy(s)}" for s in final if tidy(s)]

        # If we still have fewer than requested, pad with next best sentences
        if len(bullets) < num_points:
            remaining = [s for _, _, s in scored if s not in final]
            for s in remaining:
                if len(bullets) >= num_points:
                    break
                t = tidy(s)
                if t:
                    bullets.append(f"• {t}")

        bullets = bullets[:num_points]

        # Save to outputs/important_points.txt
        os.makedirs("outputs", exist_ok=True)
        out_path = os.path.join("outputs", "important_points.txt")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write("\n".join(bullets))

        return out_path

    except Exception as e:
        raise Exception(f"Error extracting important points: {e}")



def extract_images_from_pdf(pdf_path, output_dir):
    try:
        os.makedirs(output_dir, exist_ok=True)
        pdf_document = fitz.open(pdf_path)
        image_paths = []

        for i in range(len(pdf_document)):
            page = pdf_document[i]
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_filename = os.path.join(output_dir, f"page_{i+1}_img_{img_index+1}.png")
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)
                image_paths.append(image_filename)
        
        # Create a ZIP file of the images
        zip_filename = os.path.join(output_dir, "extracted_images.zip")
        with zipfile.ZipFile(zip_filename, 'w') as zipf:
            for img_path in image_paths:
                zipf.write(img_path, os.path.basename(img_path))  # Add image to zip

        return zip_filename  # Return the path to the ZIP file
    except Exception as e:
        raise e
def generate_pdf_presentation(pdf_path, output_dir='outputs'):
    """
    Converts a PDF into a natural-looking PowerPoint presentation with
    light backgrounds, keyword-based slide titles, and clean formatting.
    """
    try:
        import os
        import fitz  # PyMuPDF
        import pdfplumber
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from collections import Counter
        from random import choice
        import re

        os.makedirs(output_dir, exist_ok=True)

        # 🌿 Light, pastel color palette
        color_palette = [
            RGBColor(224, 242, 241),  # Light Teal
            RGBColor(255, 249, 196),  # Light Yellow
            RGBColor(227, 242, 253),  # Light Blue
            RGBColor(248, 187, 208),  # Light Pink
            RGBColor(232, 245, 233),  # Light Green
        ]

        # Utility: clean text
        def clean_text(text):
            text = text.replace('\n', ' ').replace('\t', ' ')
            return ' '.join(text.split()).strip()

        # Extract text from PDF
        text_per_page = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = clean_text(page.extract_text() or "")
                text_per_page.append(page_text)

        # Extract images per page
        pdf_document = fitz.open(pdf_path)
        image_paths_per_page = []
        for i, page in enumerate(pdf_document):
            page_images = []
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_filename = os.path.join(output_dir, f"page_{i+1}_img_{img_index+1}.png")
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)
                page_images.append(image_filename)
            image_paths_per_page.append(page_images)

        # Helper: extract a one-word heading (main keyword)
        def extract_heading(text):
            words = re.findall(r'\b[A-Z][a-zA-Z]+\b', text)  # capitalized words
            if not words:
                words = re.findall(r'\b[a-zA-Z]{5,}\b', text)  # fallback to long words
            common = Counter(words).most_common(1)
            return common[0][0].capitalize() if common else "Overview"

        # Create PowerPoint
        prs = Presentation()

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "PDF Presentation"
        slide.placeholders[1].text = "Auto-generated PowerPoint"
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(227, 242, 253)  # very light blue

        # Add slides for each PDF page
        for i, page_text in enumerate(text_per_page, start=1):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            bg_color = choice(color_palette)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg_color

            # Heading (keyword)
            heading = extract_heading(page_text)

            # Title box
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = heading
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Text content
            if page_text.strip():
                text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(3.5))
                tf = text_box.text_frame
                tf.word_wrap = True
                tf.text = page_text[:1000] + ("..." if len(page_text) > 1000 else "")
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(16)
                        run.font.color.rgb = RGBColor(60, 60, 60)

            # Add image (if available)
            if image_paths_per_page[i-1]:
                img_path = image_paths_per_page[i-1][0]
                slide.shapes.add_picture(img_path, Inches(1.5), Inches(5), width=Inches(7), height=Inches(3))

        # Outro Slide
        outro_layout = prs.slide_layouts[6]
        outro_slide = prs.slides.add_slide(outro_layout)
        outro_slide.background.fill.solid()
        outro_slide.background.fill.fore_color.rgb = RGBColor(232, 245, 233)
        box = outro_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
        tf = box.text_frame
        tf.text = "Presentation Generated Successfully!"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(28)
                run.font.color.rgb = RGBColor(33, 33, 33)

        # Save final PPT
        output_pptx = os.path.join(output_dir, "creative_pdf_presentation.pptx")
        prs.save(output_pptx)
        return output_pptx

    except Exception as e:
        raise Exception(f"Error generating presentation: {e}")





# Function to extract Q&A from PDF
def extract_qa_from_pdf(pdf_path, max_qas=10):
    """
    Extracts natural and creative Q&A pairs from a PDF file quickly.
    Uses sentence structure and keywords to generate human-like questions.
    """
    try:
        import os, re, pdfplumber
        from collections import Counter

        # Step 1: Extract text from PDF
        text = ""
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + " "

        text = re.sub(r'\s+', ' ', text.strip())
        if not text:
            raise ValueError("No text found in the PDF.")

        # Step 2: Split into sentences
        sentences = re.split(r'(?<=[.!?])\s+', text)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 30]
        if not sentences:
            raise ValueError("Not enough sentences to generate Q&A.")

        # Step 3: Find keywords for creativity
        words = re.findall(r'\b[A-Za-z]{4,}\b', text.lower())
        common_words = {w for w, _ in Counter(words).most_common(60)}

        def generate_question(sentence):
            s = sentence.strip()
            s_lower = s.lower()

            # Creative question generation
            if "because" in s_lower or "due to" in s_lower:
                return "Why does this happen?"
            elif "use" in s_lower or "used for" in s_lower:
                return "What is it used for?"
            elif "example" in s_lower:
                return "Can you give an example?"
            elif "important" in s_lower or "significant" in s_lower:
                return "Why is this important?"
            elif "result" in s_lower or "cause" in s_lower:
                return "What is the result or cause?"
            else:
                # Pick a main keyword
                keywords = [w for w in re.findall(r'\b[A-Za-z]{4,}\b', s) if w.lower() in common_words]
                if keywords:
                    key = keywords[0].capitalize()
                    return f"What is {key}?"
                else:
                    return "What does this mean?"

        # Step 4: Select meaningful sentences
        selected_sentences = []
        seen = set()
        for s in sentences:
            if len(selected_sentences) >= max_qas:
                break
            main = s.split(" ")[0:8]
            if " ".join(main) not in seen:
                selected_sentences.append(s)
                seen.add(" ".join(main))

        # Step 5: Generate Q&A pairs
        qa_pairs = []
        for s in selected_sentences:
            q = generate_question(s)
            a = s
            qa_pairs.append((q, a))

        # Step 6: Save output
        os.makedirs("outputs", exist_ok=True)
        output_file = os.path.join("outputs", "qa_output.txt")

        with open(output_file, "w", encoding="utf-8") as f:
            for i, (q, a) in enumerate(qa_pairs, 1):
                f.write(f"{i}. Q: {q}\nA: {a}\n\n")

        print(f"✅ Q&A extracted successfully: {output_file}")
        return output_file

    except Exception as e:
        raise Exception(f"Error extracting Q&A: {e}")




def merge_pdfs_with_features(pdf1, pdf2, order):
    """
    Merges two PDFs with specific features:
    1. Removes duplicate content.
    2. Maintains headings, subheadings, and structure.
    3. Allows customizable merging order.

    Args:
        pdf1 (file): First PDF file.
        pdf2 (file): Second PDF file.
        order (str): '1' for pdf1 first, '2' for pdf2 first.

    Returns:
        str: Path to the merged PDF file.
    """
    def extract_text(pdf_file):
        """Extract text content from a PDF file."""
        text = ""
        with pdfplumber.open(pdf_file) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text

    # Extract unique content
    content_pdf1 = extract_text(pdf1)
    content_pdf2 = extract_text(pdf2)
    unique_content_pdf2 = "\n".join(
        line for line in content_pdf2.splitlines() if line.strip() and line not in content_pdf1
    )

    pdf_writer = PyPDF2.PdfWriter()

    def add_pages_to_writer(pdf_file, writer):
        reader = PyPDF2.PdfReader(pdf_file)
        for page in reader.pages:
            writer.add_page(page)

    # Merge based on order
    if order == "1":
        add_pages_to_writer(pdf1, pdf_writer)
        add_pages_to_writer(pdf2, pdf_writer)  # Add all pages from PDF2
    elif order == "2":
        add_pages_to_writer(pdf2, pdf_writer)
        add_pages_to_writer(pdf1, pdf_writer)

    output_path = "merged_output.pdf"
    with open(output_path, "wb") as output_file:
        pdf_writer.write(output_file)

    return output_path


# Flask Routes
@app.route('/', methods=["GET"])
def home():
    return render_template('/index.html')


@app.route('/process-pdf', methods=['POST'])
def process_pdf():
    """
    Flask route to process and download the extracted text from an uploaded PDF.
    """
    if 'file' not in request.files:
        return jsonify({"error": "No file provided."}), 400

    pdf_file = request.files['file']
    if pdf_file.filename == '':
        return jsonify({"error": "No selected file."}), 400

    try:
        # Save the uploaded file temporarily
        temp_pdf_path = os.path.join(OUTPUT_DIR, pdf_file.filename)
        pdf_file.save(temp_pdf_path)

        # Extract and process text
        processed_text = extract_text_from_pdf(temp_pdf_path)

        # Save the processed text to a file
        output_file = os.path.join(OUTPUT_DIR, "extracted_text.txt")
        with open(output_file, "w", encoding="utf-8") as file:
            file.write(processed_text)

        # Clean up the temporary PDF file
        os.remove(temp_pdf_path)

        # Send the processed text file for download
        return send_file(output_file, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/process-tables', methods=['POST'])
def process_tables():
    try:
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Extract tables from PDF
        tables = extract_tables_from_pdf(pdf_path)
        if not tables:
            return jsonify({"error": "No tables found in the PDF."}), 404

        # Save tables to an Excel file
        output_path = os.path.join(OUTPUT_DIR, f"tables_{file.filename}.xlsx")
        with pd.ExcelWriter(output_path) as writer:
            for i, table in enumerate(tables):
                df = pd.DataFrame(table)
                df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)

        # Send the Excel file as an attachment
        return send_file(output_path, as_attachment=True, download_name=f"tables_{file.filename}.xlsx", mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/summarize', methods=['POST'])
def summarize_pdf():
    try:
        # Get the uploaded file
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file provided"}), 400

        # Save the file in the UPLOAD_DIR
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Extract text from the uploaded PDF
        text = extract_text_from_pdf(pdf_path)
        text = preprocess_text(text)

        # Chunk the text and summarize
        text_chunks = chunk_text_by_tokens(text)
        summaries = []
        for chunk in text_chunks:
            input_ids = tokenizer(chunk, return_tensors="pt", truncation=True).input_ids
            summary_ids = model.generate(input_ids, max_length=512, min_length=50, no_repeat_ngram_size=3)
            summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
            summaries.append(summary)

        summary_text = " ".join(summaries)

        # Save the summary to a file in the UPLOAD_DIR
        summary_file_path = os.path.join(UPLOAD_DIR, "summary.txt")
        with open(summary_file_path, "w", encoding="utf-8") as summary_file:
            summary_file.write(summary_text)

        # Automatically download the file
        return send_file(summary_file_path, as_attachment=True, download_name="summary.txt")

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/extract-images', methods=['POST'])
def extract_images():
    try:
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        output_dir = os.path.join(OUTPUT_DIR, f"images_{file.filename}")
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        zip_file_path = extract_images_from_pdf(pdf_path, output_dir)
        if not zip_file_path:
            return jsonify({"error": "No images found in the PDF."}), 404

        return send_file(zip_file_path, as_attachment=True, mimetype='application/zip', download_name="extracted_images.zip")

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/extract-points', methods=['POST'])
def extract_points():
    """
    Flask route to upload a PDF, extract important points, and download the result.
    """
    try:
        # Step 1: Validate the uploaded file
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        # Step 2: Save the uploaded PDF
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Step 3: Extract text from the PDF
        text = extract_text_from_pdf(pdf_path)
        if not text:
            return jsonify({"error": "No text found in the PDF."}), 404

        # Step 4: Extract important points from the text
        output_path = extract_important_points(text)

        # Step 5: Return the generated text file for download
        return send_file(output_path, as_attachment=True, mimetype='text/plain', download_name="important_points.txt")

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/upload-pdf', methods=['POST'])
def upload_pdf():
    try:
        # Check if a file is uploaded
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        # Save the uploaded PDF
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Generate the PowerPoint presentation
        pptx_path = generate_pdf_presentation(pdf_path)

        # Send the generated PowerPoint presentation for download
        return send_file(pptx_path, as_attachment=True, download_name="pdf_presentation.pptx")

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/extract-q&a', methods=['POST'])
def upload_and_extract_qa():
    try:
        # Check if a file is uploaded
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if not file.filename.endswith('.pdf'):
            return jsonify({"error": "Only PDF files are allowed"}), 400

        # Save the uploaded file temporarily
        pdf_path = file.filename
        file.save(pdf_path)

        # Extract Q&A and get the output file
        output_file = extract_qa_from_pdf(pdf_path)

        # Send the file as a downloadable response
        if output_file:
            return send_file(output_file, as_attachment=True)
        else:
            return jsonify({"error": "Failed to extract Q&A"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500



@app.route('/merge', methods=['POST'])
def merge_pdfs():
    if 'pdf1' not in request.files or 'pdf2' not in request.files:
        return jsonify({"error": "Both PDF files are required."}), 400

    pdf1 = request.files['pdf1']
    pdf2 = request.files['pdf2']
    order = request.form.get('order', '1')

    try:
        merged_pdf_path = merge_pdfs_with_features(pdf1, pdf2, order)
        return send_file(merged_pdf_path, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    '''finally:
        if os.path.exists(merged_pdf_path):
            os.remove(merged_pdf_path)'''



# Cleanup files periodically (optional)
def cleanup():
    for folder in [UPLOAD_DIR, OUTPUT_DIR]:
        for filename in os.listdir(folder):
            file_path = os.path.join(folder, filename)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            except Exception as e:
                print(f"Failed to delete {file_path}. Reason: {e}")

# Main Driver
if __name__ == '__main__':
    app.run(debug=False, port=5000)




